from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


def add_bullets(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(24)


def add_code_slide(prs: Presentation, title: str, code_text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(12.3), Inches(5.4))
    tf = textbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.name = "Courier New"
    p.font.size = Pt(14)


def main() -> None:
    prs = Presentation()

    # 1) Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Raspberry Pi Pothole Detection"
    slide.placeholders[1].text = "ELE408 Course Project\nNoah Varghese"

    # 2) Background
    add_bullets(
        prs,
        "Background",
        [
            "Potholes damage vehicles and create safety hazards.",
            "Dashcam road imagery enables visual pothole detection.",
            "Edge computing on Raspberry Pi supports low-latency, offline operation.",
            "Goal: build a practical in-vehicle detector using camera + GPS + local alerts.",
        ],
    )

    # 3) Problem Statement
    add_bullets(
        prs,
        "Problem Statement",
        [
            "Detect potholes in real road scenes with an onboard camera.",
            "Store pothole locations only when confidence is high.",
            "Ignore non-pothole frames to avoid unnecessary storage.",
            "Warn drivers when approaching known potholes in real time.",
        ],
    )

    # 4) Basic Idea
    add_bullets(
        prs,
        "Basic Idea",
        [
            "Capture one image every X seconds from Arducam.",
            "Run YOLOv26n-cls inference on Raspberry Pi 5.",
            "If predicted pothole: read NEO-6M GPS fix and store in SQLite.",
            "Continuously compute distance to nearest stored pothole for alerts.",
        ],
    )

    # 5) Implementation
    add_bullets(
        prs,
        "Implementation",
        [
            "Hardware: Raspberry Pi 5 + Arducam 5MP + u-blox NEO-6M (UART).",
            "Software: Python, Ultralytics YOLO, OpenCV, gpsd-py3, SQLite.",
            "Services: capture, inference, GPS, database, proximity, notifier.",
            "Main loop coordinates capture -> classify -> geotag -> dedup -> alert.",
        ],
    )

    # 6) Code
    code_text = """# main.py (runtime loop excerpt)
image_path = capture.capture()
fix = gps.get_fix()
prediction = infer.predict(image_path)

if prediction["is_pothole"] and fix is not None:
    db.insert_pothole(
        latitude=fix.latitude,
        longitude=fix.longitude,
        detected_at=fix.timestamp,
        confidence=float(prediction["confidence"]),
        image_id_optional=image_path.name,
    )

nearest = db.nearest_distance_m(fix.latitude, fix.longitude) if fix else None
if proximity.should_alert(nearest):
    notifier.notify(f"Pothole ahead in {nearest:.1f} meters.")
"""
    add_code_slide(prs, "Code", code_text)

    # 7) Evaluation
    add_bullets(
        prs,
        "Evaluation",
        [
            "Measured metrics: avg inference latency, frames/min, GPS fix usage, potholes saved.",
            "Log parser script extracts runtime summary from detector logs.",
            "Dedup distance reduces repeated inserts for the same pothole pass.",
            "Next step: collect route-based precision/recall from field trials.",
        ],
    )

    # 8) Video Documentation
    add_bullets(
        prs,
        "Video Documentation",
        [
            "Demo plan: run detector in moving vehicle with live console output.",
            "Show: image capture cadence, inference labels/confidence, and DB inserts.",
            "Show: proximity alerts triggering near known pothole coordinates.",
            "Insert your final demo link(s) here before presentation.",
        ],
    )

    # 9) References
    add_bullets(
        prs,
        "References",
        [
            "Ultralytics YOLO documentation: https://docs.ultralytics.com",
            "Raspberry Pi documentation: https://www.raspberrypi.com/documentation/",
            "gpsd project: https://gpsd.io/",
            "u-blox NEO-6M product documentation and integration guides.",
        ],
    )

    output = Path("Pothole_Detection_Presentation.pptx")
    prs.save(output)
    print(f"Created: {output.resolve()}")


if __name__ == "__main__":
    main()
